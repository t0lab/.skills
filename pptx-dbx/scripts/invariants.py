#!/usr/bin/env python3
"""Invariants ruler for .pptx decks.

Pure numbers in, pure numbers out — no PASS/FAIL, no thresholds. The reviewer
applies thresholds against spec.md. This tool exists because bbox overlaps and
header/footer band clearance are structural facts that measure.py's
per-shape-pair queries cannot surface at a glance.

Outputs per slide:
  - header_band_bottom / footer_band_top (from Header_*_Bar / Footer_*_Bar)
  - For every content shape: clearance_top, clearance_bottom (EMU)
  - Every pair of non-same-group shapes whose bboxes intersect with area > 0

Same-group detection: two shapes share a group if their objectName prefix
(everything before the last underscore) matches.
  CharCard_Perception_Container ~ CharCard_Perception_Title  → same group
  Arch_Arrow2_Line              ~ Arch_Eyebrow_Label          → different
"""

import argparse
import sys
from pptx import Presentation


LINE_INFLATE_EMU = 12700  # ~1pt; inflates zero-dim shapes (lines) so their strokes count for overlap


def bbox(shape):
    l, t = shape.left, shape.top
    r, b = l + shape.width, t + shape.height
    # Zero-dim shapes (lines drawn with h=0 or w=0) need inflation so stroke-area overlap is detectable
    if shape.width == 0:
        l -= LINE_INFLATE_EMU
        r += LINE_INFLATE_EMU
    if shape.height == 0:
        t -= LINE_INFLATE_EMU
        b += LINE_INFLATE_EMU
    return (l, t, r, b)


def intersect_area(a, b):
    al, at, ar, ab = a
    bl, bt, br, bb = b
    iw = max(0, min(ar, br) - max(al, bl))
    ih = max(0, min(ab, bb) - max(at, bt))
    return iw * ih


def same_group(a, b):
    pa = "_".join(a.split("_")[:-1])
    pb = "_".join(b.split("_")[:-1])
    return bool(pa) and pa == pb


def fully_inside(inner, outer):
    il, it, ir, ib = inner
    ol, ot, or_, ob = outer
    return il >= ol and it >= ot and ir <= or_ and ib <= ob


def is_container(name):
    return name.endswith("_Container")


def collect(shapes, out):
    for s in shapes:
        if s.shape_type == 6:  # GROUP
            collect(s.shapes, out)
        else:
            if s.has_text_frame or s.shape_type is not None:
                try:
                    out.append((s.name, bbox(s)))
                except Exception:
                    pass


def is_chrome(name):
    return name.startswith("Header_") or name.startswith("Footer_")


def is_background(name):
    # Full-slide background shapes cover the canvas by design — not an overlap bug.
    return "_Background_" in name or name.endswith("_Background") or name.startswith("Background_")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--slide", type=int, default=None)
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    print(f"# Invariants report: {args.pptx}")
    print(f"# slide_dims_emu: {prs.slide_width} x {prs.slide_height}")

    for i, slide in enumerate(prs.slides, 1):
        if args.slide and i != args.slide:
            continue
        shapes = []
        collect(slide.shapes, shapes)

        print(f"\n## Slide {i}")

        header_bottom = None
        footer_top = None
        for name, (l, t, r, b) in shapes:
            if name.startswith("Header_") and name.endswith("_Bar"):
                header_bottom = b if header_bottom is None else max(header_bottom, b)
            if name.startswith("Footer_") and name.endswith("_Bar"):
                footer_top = t if footer_top is None else min(footer_top, t)
        print(f"header_band_bottom_emu: {header_bottom}")
        print(f"footer_band_top_emu: {footer_top}")

        if header_bottom is not None or footer_top is not None:
            print("\n### Content-shape clearance (EMU)")
            print("# clearance_top = shape.top - header_band_bottom")
            print("# clearance_bottom = footer_band_top - shape.bottom")
            for name, (l, t, r, b) in shapes:
                if is_chrome(name):
                    continue
                ct = (t - header_bottom) if header_bottom is not None else None
                cb = (footer_top - b) if footer_top is not None else None
                print(f"  {name}: clearance_top={ct}, clearance_bottom={cb}")

        print("\n### Overlapping shape pairs (bbox intersection area > 0)")
        print("# same-group pairs suppressed (share objectName prefix)")
        found = 0
        for j in range(len(shapes)):
            for k in range(j + 1, len(shapes)):
                na, ba = shapes[j]
                nb, bbo = shapes[k]
                if same_group(na, nb):
                    continue
                if is_chrome(na) and is_chrome(nb):
                    continue
                if is_background(na) or is_background(nb):
                    continue
                # Parent-child: a child shape fully inside a *_Container is not an overlap bug
                if is_container(na) and fully_inside(bbo, ba):
                    continue
                if is_container(nb) and fully_inside(ba, bbo):
                    continue
                area = intersect_area(ba, bbo)
                if area > 0:
                    print(f"  {na} ∩ {nb}: area_emu2={area}")
                    found += 1
        if found == 0:
            print("  (none)")


if __name__ == "__main__":
    main()
